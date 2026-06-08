from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE = RGBColor(0x0D, 0x3B, 0x6E)
MED_BLUE = RGBColor(0x1A, 0x6E, 0xA8)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC8, 0x96, 0x2E)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT,
                 font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_slide(slide, title_text, bullets):
    _set_slide_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 title_text, font_size=28, bold=True, color=DARK_BLUE)
    # Accent line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.0), Inches(2.5), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    _add_textbox(slide, Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.5),
                 bullets, font_size=18, color=DARK_TEXT)


def generate_app_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════
    # Slide 1 – Title
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide, DARK_BLUE)

    _add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
                 'منظومة توثيق أراضي طرح النهر',
                 font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(2.8), Inches(8), Inches(0.6),
                 'وزارة الموارد المائية والري — جمهورية مصر العربية',
                 font_size=20, color=GOLD, alignment=PP_ALIGN.CENTER)

    # Separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.7), Inches(3), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = GOLD
    sep.line.fill.background()

    _add_textbox(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
                 'نظام إلكتروني متكامل لتوثيق ورصد التعديات على أراضي طرح النهر',
                 font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(6.5), Inches(8), Inches(0.4),
                 f'الإصدار 1.0 | {datetime.now().year}',
                 font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════
    # Slide 2 – Introduction / Problem
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'مقدمة', """\
تواجه أراضي طرح نهر النيل تحديات كبيرة بسبب التعديات والاستغلال غير القانوني، مما يؤثر على:
• حق المواطنين في الانتفاع بأراضي الدولة
• المجرى المائي للنيل وصيانة الجسور
• الإيرادات المستحقة للدولة

تأتي المنظومة كحل رقمي متكامل لتوثيق ورصد وإدارة هذه التعديات بشكل منهجي وشفاف.""")
    # ══════════════════════════════════════════
    # Slide 3 – Overview
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'عن المنظومة', """\
منظومة توثيق أراضي طرح النهر هي تطبيق ويب شامل يهدف إلى:
• تسجيل التعديات على أراضي طرح النهر ببيانات دقيقة (إحداثيات، مساحات، خرائط مساحية)
• ربط التعديات بالجهات الإدارية (محافظة، مركز، قرية) وببيانات الحوض
• توثيق الصور الحقلية والخرائط المساحية لكل تعدٍّ
• تتبع فترات الاستغلال المتعددة للأرض وحساب مقابل الانتفاع
• مقارنة صور الأقمار الصناعية عبر السنوات لرصد التغيرات
• إصدار تقارير رسمية (Excel، PDF) للإجراءات القانونية والإدارية""")

    # ══════════════════════════════════════════
    # Slide 4 – Interactive Map
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'الخريطة التفاعلية', """\
• خريطة ويب تفاعلية تعرض التعديات على خريطة الأساس (Esri Satellite / OpenStreetMap)
• عرض المضلعات المساحية (Polygons) المستوردة من ملفات Shapefile
• تجميع العلامات (Marker Clustering) لتسهيل التصفح
• لوحة جانبية للبحث والتصفية حسب: المحافظة، المركز، الحالة، المساحة
• عرض تفاصيل كل تعدٍّ مع الصور والخرائط المساحية بصيغة PDF
• طبقة صور الأقمار الصناعية (Sentinel Hub WMS) للمقارنة المباشرة""")

    # ══════════════════════════════════════════
    # Slide 5 – Satellite Imagery
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'صور الأقمار الصناعية والرصد الفضائي', """\
• الاتصال بخدمة Sentinel Hub Process API لجلب صور Sentinel-2 L2A
• مقارنة صور سنوات متعددة لاكتشاف التغيرات (Before/After)
• كشف التغيير على مستوى البكسل: حساب متوسط الفرق، نسبة التغيير، المساحة التقديرية
• إنشاء خرائط حرارية (Heatmap) بطبقة حمراء فوق المناطق المتغيرة
• تخزين الصور والخرائط الحرارية في قاعدة البيانات للتسريع
• تقارير PDF فضائية كاملة بالصور والإحصائيات""")

    # ══════════════════════════════════════════
    # Slide 6 – Data Entry & Shapefile Import
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'إدخال البيانات واستيراد Shapefile', """\
إدخال يدوي:
• تسجيل التعديات بإحداثيات دقيقة (Lat/Lng) ومساحات متعددة
• إرفاق الصور الحقلية والخرائط المساحية (PDF)
• ربط التعديات بالجهات الإدارية وأنواع الاستغلال

استيراد Shapefile:
• رفع ملفات ESRI Shapefile (ZIP) مع إعادة الإسقاط إلى WGS84
• تعيين الحقول تلقائياً (عربي/إنجليزي) مع إمكانية التعديل اليدوي
• استيراد إلى 10,000 سجل مع المضلعات الهندسية (GeoJSON)
• استخراج Pcode المحافظة من بيانات الشكل""")

    # ══════════════════════════════════════════
    # Slide 7 – Usage Tracking
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'تتبع الاستغلال ومقابل الانتفاع', """\
• تسجيل فترات استغلال متعددة لكل قطعة أرض مع بيانات المنتفع
• حساب مقابل الانتفاع استناداً إلى القرارات الوزارية التاريخية (1987–2026)
• تسعير حسب المناطق (الوراق، أسوان/الأقصر/أسيوط، محافظات النيل الأخرى)
• معامل تخفيض: 1.0 للمخالفات (قرار 149) و 0.5 للتراخيص (قرار 148)
• تطبيق تلقائي للقرار الوزاري المناسب حسب الفترة الزمنية
• طباعة سجل استغلال واحد أو全て السجلات مع تفاصيل الحساب""")

    # ══════════════════════════════════════════
    # Slide 8 – Reports & Exports
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'التقارير والتصدير', """\
• تقارير Excel: جدول بيانات RTL كامل بتنسيق احترافي وألوان حسب الحالة
• تقارير PDF: تنسيق أفقي A4 مع جداول محترمة وترتيب RTL
• تقارير الأقمار الصناعية: صور متعددة السنوات + خرائط حرارية + إحصائيات
• دعم كامل للغة العربية (تشكيل الأحرف وترتيب النص)
• سجل تدقيق كامل (Audit Log) لجميع العمليات مع إمكانية التصدير""")

    # ══════════════════════════════════════════
    # Slide 9 – Chatbot
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'المساعد الذكي (Chatbot)', """\
• chatbot باللغة العربية للاستعلام عن البيانات
• دعم أنماط الأسئلة: عدد، عرض، مجموع، متوسط، تفاصيل
• استخراج الكيانات: المنطقة، القرار (148/149)، الحالة، المحافظة
• أمثلة: «كم عدد التعديات في محافظة الجيزة؟»
          «اعرض التعديات بقرار 149 في محافظة أسوان»
          «إجمالي المساحات المتعدى عليها في القليوبية»
• تكامل مباشر مع قاعدة البيانات""")
    # ══════════════════════════════════════════
    # Slide 10 – Admin Dashboard
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'لوحة إدارة النظام', """\
• إحصائيات شاملة: إجمالي التعديات، قيد الانتظار، معتمد، مرفوض، إجمالي المساحات
• إدارة المستخدمين: إنشاء، تعديل، تعطيل، تعيين الأدوار والمحافظات
• سجل التدقيق: عرض وتصفية وتصدير جميع العمليات
• إدارة القرارات الوزارية: إضافة، رفع PDF، استيراد أسعار Excel
• إدارة أسعار الاستخدام حسب القرار والمنطقة ونوع الاستغلال""")

    # ══════════════════════════════════════════
    # Slide 11 – User Roles
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'صلاحيات المستخدمين', """\
المشاهد (Viewer):
• عرض الخريطة والسجلات المعتمدة وتوليد التقارير

مدخل البيانات (Data Entry):
• إضافة وتعديل التعديلات، رفع الصور والخرائط المساحية، تسجيل فترات الاستغلال

المشرف (Supervisor):
• اعتماد ورفض التعديلات وسجلات الاستغلال، استيراد Shapefile

مدير النظام (Manager):
• صلاحية كاملة: لوحة الإدارة، إدارة المستخدمين، سجل التدقيق، إدارة القرارات الوزارية""")

    # ══════════════════════════════════════════
    # Slide 12 – Technical Architecture
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'التقنيات المستخدمة', """\
الخلفية (Backend):
• Django 4.2 (Python) – إطار العمل الأساسي
• PostgreSQL مع PostGIS – قاعدة البيانات
• GeoPandas + Fiona + Shapely – معالجة البيانات المكانية
• OpenPyXL + ReportLab – تصدير Excel و PDF

الواجهة (Frontend):
• Leaflet.js – الخريطة التفاعلية
• MarkerCluster – تجميع العلامات
• Font Awesome 6 – الأيقونات
• PDF.js – عرض الخرائط المساحية

الاستشعار عن بُعد:
• Sentinel Hub Process API – صور Sentinel-2
• NumPy + Pillow – معالجة الصور وكشف التغيير

النشر:
• Gunicorn + WhiteNoise – خادم الإنتاج
• Render – منصة الاستضافة""")

    # ══════════════════════════════════════════
    # Slide 13 – Benefits
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bullet_slide(slide, 'فوائد المنظومة', """\
• توثيق منهجي وشفاف لجميع التعديات على أراضي طرح النهر
• أرشفة رقمية كاملة (بيانات، صور، خرائط، مستندات)
• تحليل زمني باستخدام صور الأقمار الصناعية لاكتشاف التعديات الجديدة
• حوكمة الإجراءات عبر سير عمل اعتماد متعدد المستويات
• سجل تدقيق كامل لضمان المساءلة والشفافية
• حساب دقيق لمقابل الانتفاع وفقاً للقرارات الوزارية
• توفير الوقت والجهد مقارنة بالنظام الورقي التقليدي""")

    # ══════════════════════════════════════════
    # Slide 14 – Thank You
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BLUE)

    _add_textbox(slide, Inches(1), Inches(2.0), Inches(8), Inches(1.0),
                 'والله ولي التوفيق',
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.6),
                 'وزارة الموارد المائية والري',
                 font_size=20, color=GOLD, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.5),
                 'منظومة توثيق أراضي طرح النهر',
                 font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(5.0), Inches(8), Inches(0.5),
                 'جميع الحقوق محفوظة © وزارة الموارد المائية والري',
                 font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    return prs
